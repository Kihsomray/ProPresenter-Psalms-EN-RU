# coding: utf8

# Author: Michael Yarmoshik
# Version: 1.0.0 (2024-02-03)
# https://github.com/Kihsomray/Psalm-Offsets-EN-RU

# This project is not perfect by any means. However, it allowed for the streamlining
# fixing the offsets between the English (ESV) and Russian (SYNO) translations.

# Each query requires the English location, the English chapter that includes verses,
# and the Russian chapter that includes verses.

# Essentially, it creates a PowerPoint that can be imported into ProPresenter, where
# a theme can easily be applied to make it visually appealing.

from pptx import Presentation
from pptx.util import Inches
import re


# What is the English chapter are we considering?
english_chapter = 150

# What is the English text?
english = """
1 Praise the Lord!
Praise God in his sanctuary;
praise him in his mighty heavens!
2Praise him for his mighty deeds;
praise him according to his excellent greatness!
3Praise him with trumpet sound;
praise him with lute and harp!
4Praise him with tambourine and dance;
praise him with strings and pipe!
5Praise him with sounding cymbals;
praise him with loud clashing cymbals!
6Let everything that has breath praise the Lord!
Praise the Lord!
"""

# What is the Russian text?
russian = """1 Хвалите Бога во святыне Его, хвалите Его на тверди силы Его. 2Хвалите Его по могуществу Его, хвалите Его по множеству величия Его. 3 Хвалите Его со звуком трубным, хвалите Его на псалтири и гуслях. 4 Хвалите Его с тимпаном и ликами, хвалите Его на струнах и органе. 5Хвалите Его на звучных кимвалах, хвалите Его на кимвалах громогласных. 6 Все дышащее да хвалит Господа! Аллилуия.
"""

offsets = [
    ["1:1-2:12", "X:X"],
    ["3:1-9:20", "X:+1"],
    ["10:1-1", "9:22"],
    ["10:2-18", "9:21"],
    ["11:1-7", "10:X"],
    ["12:1-13:4", "-1:+1"],
    ["13:5-6", "12:6"],  # 5 --> 6a, 6 --> 6b
    ["14:1-17:15", "-1:X"],
    ["18:1-22:31", "-1:+1"],
    ["23:1-29:11", "-1:X"],
    ["30:1-31:24", "-1:+1"],
    ["32:1-33:22", "-1:X"],
    ["34:1-22", "33:+1"],
    ["35:1-28", "34:X"],
    ["36:1-12", "35:+1"],
    ["37:1-40", "36:X"],
    ["38:1-42:11", "-1:+1"],
    ["43:1-5", "42:X"],
    ["44:1-49:20", "-1:+1"],
    ["50:1-23", "49:X"],
    ["51:1-52:9", "-1:+2"],
    ["53:1-6", "52:+1"],
    ["54:1-7", "53:+2"],
    ["55:1-59:17", "-1:+1"],
    ["60:1-12", "59:+2"],
    ["61:1-65:13", "-1:+1"],
    ["66:1-20", "65:X"],
    ["67:1-70:5", "-1:+1"],
    ["71:1-74:23", "-1:X"],
    ["75:1-77:20", "-1:+1"],
    ["78:1-79:13", "-1:X"],
    ["80:1-81:16", "-1:+1"],
    ["82:1-8", "81:X"],
    ["83:1-85:13", "-1:+1"],
    ["86:1-17", "85:X"],
    ["87:1-7", "86:X"],  # 2 --> 2a, 3, --> 2b
    ["88:1-90:5", "-1:+1"],  # 5 --> 6a, 6 --> 6b
    ["90:6-91:16", "-1:X"],
    ["92:1-15", "91:+1"],
    ["93:1-101:8", "-1:X"],
    ["102:1-28", "101:+1"],
    ["103:1-107:43", "-1:X"],
    ["108:1-13", "107:+1"],
    ["109:1-114:8", "-1:X"],
    ["115:1-18", "113:+8"],
    ["116:1-9", "114:X"],  # 8 --> 8a, 9 --> 8b
    ["116:10-19", "115:-9"],
    ["117:1-139:24", "-1:X"],
    ["140:1-13", "-1:+1"],
    ["141:1-147:11", "-1:X"],
    ["147:12-20", "147:-11"],
    ["148:1-150:6", "X:X"]
]


def form(string):

    # Empty list
    verses_list = []

    # Basic conversion
    string = re.sub(re.compile(r"([0-9]) "), r"\1", string.replace("\n", " "))

    # Use a regular expression to match verse numbers and text
    pattern = re.compile(r'(\d+)([^\d]+)')

    # Find all matches in the text
    matches = re.findall(pattern, string)

    # Create a list of lists containing verse number and verse text
    for match in matches:
        verses_list.insert(int(match[0]), match[1].strip())

    return verses_list


def connect_verses(english_verses, russian_verses):
    connected_verses = []

    for offset in offsets:
        en = offset[0].split("-")
        ru = offset[1].split(":")[1].replace("X", "0")

        en_c_start = int(en[0].split(":")[0])
        en_v_start = int(en[0].split(":")[1])

        en_c_end = en_c_start
        en_v_end = en_v_start

        if len(en[1].split(":")) == 2:
            en_c_end = int(en[1].split(":")[0])
            en_v_end = int(en[1].split(":")[1])
        else:
            en_v_end = int(en[1])

        if en_c_start <= english_chapter <= en_c_end:

            # get the russian chapter
            russian_chapter = int(offset[1].split(":")[1].replace("X", "0"))
            if russian_chapter < 2:
                russian_chapter += english_chapter

            for i in range(len(english_verses)):

                # within range
                if en_v_start <= i + 1 and not (english_chapter == en_c_end and i + 1 > en_v_end):
                    ru_v = int(ru) + i

                    connected_verses.append([
                        english_verses[i],
                        english_chapter,
                        i + 1,
                        russian_verses[ru_v],
                        russian_chapter,
                        ru_v + 1
                    ])
                i += 1

    return connected_verses


connected_verses = connect_verses(form(english), form(russian))


def add_blank_slide(prs):
    slide_layout = prs.slide_layouts[6]  # Use layout 6 for a blank slide
    slide = prs.slides.add_slide(slide_layout)
    return slide


def add_text_box(slide, name, text, left, top, width, height):
    text_box = slide.shapes.add_textbox(left, top, width, height)
    p = text_box.text_frame
    p.text = text
    text_box.name = name  # Set a custom name for the shape
    return text_box


# Create a presentation
presentation = Presentation()

presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

for verse in connected_verses:
    # Add a blank slide
    slide = add_blank_slide(presentation)

    # Add text boxes with specified names
    verse1_box = add_text_box(slide, "Verse 1", verse[0], Inches(1), Inches(1), Inches(2), Inches(0.5))
    ref1_box = add_text_box(slide, "Reference 1", "Psalm " + str(verse[1]) + ":" + str(verse[2]), Inches(1), Inches(2), Inches(2), Inches(0.5))
    verse2_box = add_text_box(slide, "Verse 2", verse[3], Inches(1), Inches(3), Inches(2), Inches(0.5))
    ref2_box = add_text_box(slide, "Reference 2", "Псалтирь " + str(verse[4]) + ":" + str(verse[5]), Inches(1), Inches(4), Inches(2), Inches(0.5))


# Save the presentation
presentation.save("Psalm " + str(english_chapter) + ".pptx")
